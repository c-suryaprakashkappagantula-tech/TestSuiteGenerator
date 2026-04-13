"""
Generate TSG V4.1 Presentation Deck — using sample deck as template for logos/branding.
White background, professional colors, arrows on architecture slide.
Run: python TSG_V4_Presentation.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from datetime import datetime

# Use sample deck as template — keep all slides, clear text, add our content
prs = Presentation('sample_deck.pptx')

# We'll reuse the 9 slides from the sample — just clear and replace content
# First, let's clear all text from all shapes on all slides
for slide in prs.slides:
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            # Clear all text
            for para in shape.text_frame.paragraphs:
                para.text = ''
        # Mark non-placeholder shapes for removal (custom boxes, etc.)
        if not shape.is_placeholder and shape.shape_type != 14:  # keep placeholders
            shapes_to_remove.append(shape)
    # Remove custom shapes
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

# Colors — matching sample deck's dark background
NAVY = RGBColor(0x0B, 0x1D, 0x39)
DARK_TEXT = RGBColor(0xFF, 0xFF, 0xFF)  # white text on dark bg
GRAY_TEXT = RGBColor(0xCC, 0xCC, 0xCC)  # light gray on dark bg
LIGHT_BG = RGBColor(0x10, 0x14, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_PURPLE = RGBColor(0x6C, 0x3F, 0xC0)
ACCENT_BLUE = RGBColor(0x25, 0x63, 0xEB)
ACCENT_GREEN = RGBColor(0x16, 0xA3, 0x4A)
ACCENT_ORANGE = RGBColor(0xEA, 0x88, 0x0B)
ACCENT_CYAN = RGBColor(0x05, 0x96, 0xB5)
ACCENT_RED = RGBColor(0xDC, 0x26, 0x26)


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def add_bullets(slide, left, top, width, height, items, size=13, color=GRAY_TEXT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = '•  ' + item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(3)
    return tf


def add_stat(slide, left, top, value, label, accent):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(2.0), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = accent
    shape.line.width = Pt(2)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(value)
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = accent
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(9)
    p2.font.color.rgb = GRAY_TEXT
    p2.alignment = PP_ALIGN.CENTER


def add_box(slide, left, top, w, h, label, accent, font_size=9):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = accent
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    for li, line in enumerate(label.split('\n')):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_TEXT if li == 0 else GRAY_TEXT
        p.font.bold = (li == 0)
        p.alignment = PP_ALIGN.CENTER
    return shape

def add_arrow(slide, x1, y1, x2, y2, color=ACCENT_BLUE):
    """Add a connector arrow between two points."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(2)


# ================================================================
# SLIDE 1: Title (reuse slide 0)
# ================================================================
s = prs.slides[0]
add_text(s, 1.5, 1.5, 7, 1.0, 'TSG — Test Suite Generator V4.1', size=36, bold=True, color=WHITE)
add_text(s, 1.5, 2.5, 7, 0.6, 'AI-Enhanced Test Suite Generation Platform', size=20, color=ACCENT_PURPLE)
add_text(s, 1.5, 3.3, 7, 0.5, 'Chalk + Jira + Attachments + AI → Production-Ready Test Suites', size=14, color=GRAY_TEXT)
add_text(s, 1.5, 4.5, 7, 0.4, '%s | Charter Communications | TMO Provisioning QA' % datetime.now().strftime('%B %Y'), size=12, color=GRAY_TEXT)
add_text(s, 1.5, 5.0, 7, 0.4, 'Powered by Amazon Q AI', size=12, bold=True, color=ACCENT_GREEN)


# ================================================================
# SLIDE 2: Objective & Goal (reuse slide 1)
# ================================================================
s = prs.slides[1]
add_text(s, 0.5, 0.3, 9, 0.5, 'Objective & Goal', size=28, bold=True, color=WHITE)

add_text(s, 0.5, 1.0, 4.5, 0.3, 'OBJECTIVE', size=15, bold=True, color=ACCENT_PURPLE)
add_bullets(s, 0.5, 1.3, 4.5, 2.5, [
    'Automate test suite creation from Chalk, Jira, and design docs',
    'Generate production-ready Excel suites with structured steps and traceability',
    'Apply AI-powered gap analysis to catch what rule-based engines miss',
    'Eliminate manual test case writing — hours to minutes per feature',
])

add_text(s, 5.5, 1.0, 4.5, 0.3, 'GOAL', size=15, bold=True, color=ACCENT_GREEN)
add_bullets(s, 5.5, 1.3, 4.5, 2.5, [
    'Reduce test suite preparation time by 90%+ per feature',
    'Achieve consistent, human-quality test cases across all features',
    'Enable batch generation — multiple features in one click',
    'Provide complete AC traceability and priority-based execution guidance',
])

add_stat(s, 0.5, 4.3, '128', 'Features\nSupported', ACCENT_PURPLE)
add_stat(s, 2.9, 4.3, '70+', 'TCs per\nFeature', ACCENT_BLUE)
add_stat(s, 5.3, 4.3, '6', 'Pipeline\nBlocks', ACCENT_GREEN)
add_stat(s, 7.7, 4.3, '90%', 'Time\nSaved', ACCENT_ORANGE)


# ================================================================
# SLIDE 3: Architecture with Arrows (reuse slide 2)
# ================================================================
s = prs.slides[2]
add_text(s, 0.5, 0.3, 9, 0.5, 'Architecture — Pipeline Flow', size=28, bold=True, color=WHITE)
add_text(s, 0.5, 0.8, 9, 0.3, 'Block-based execution with self-heal retry — each block retries up to 2x on failure', size=11, color=GRAY_TEXT)

# Input sources (left column)
add_box(s, 0.3, 1.5, 2.0, 0.8, 'Chalk Pages\n(Feature Details)', ACCENT_PURPLE)
add_box(s, 0.3, 2.6, 2.0, 0.8, 'Jira\n(Stories & AC)', ACCENT_BLUE)
add_box(s, 0.3, 3.7, 2.0, 0.8, 'LLD / HLD\nDocuments', ACCENT_CYAN)

# Normalized layer (center-left)
add_box(s, 3.0, 2.3, 2.2, 1.2, 'Normalized Input\n(Structured JSON\n& DB Cache)', ACCENT_GREEN, 10)

# Engine (center-right)
add_box(s, 5.8, 2.3, 2.2, 1.2, 'Deterministic\nTest Engine\n(Rule-Based)', ACCENT_ORANGE, 10)

# AI + Output (right column)
add_box(s, 8.3, 1.5, 2.0, 1.0, 'AI Analysis\nMissing Scenarios\nCoverage Gaps', ACCENT_RED, 9)
add_box(s, 8.3, 2.8, 2.0, 0.8, 'Human Review\n& Approval', ACCENT_CYAN, 9)
add_box(s, 8.3, 3.9, 2.0, 0.8, 'Approved\nTest Suites', ACCENT_GREEN, 9)

# Arrows: inputs → normalized
add_arrow(s, 2.3, 1.9, 3.0, 2.7, ACCENT_PURPLE)
add_arrow(s, 2.3, 3.0, 3.0, 2.9, ACCENT_BLUE)
add_arrow(s, 2.3, 4.1, 3.0, 3.1, ACCENT_CYAN)
# Normalized → Engine
add_arrow(s, 5.2, 2.9, 5.8, 2.9, ACCENT_GREEN)
# Engine → AI
add_arrow(s, 8.0, 2.9, 8.3, 2.0, ACCENT_ORANGE)
# AI → Review
add_arrow(s, 9.3, 2.5, 9.3, 2.8, ACCENT_RED)
# Review → Approved
add_arrow(s, 9.3, 3.6, 9.3, 3.9, ACCENT_CYAN)

add_text(s, 0.3, 5.0, 10, 0.3, 'Tech:  Python 3.13  |  Streamlit  |  Playwright  |  SQLite  |  OpenPyXL  |  LLM-Ready (OpenAI / Bedrock / Ollama)', size=10, color=GRAY_TEXT)


# ================================================================
# SLIDE 4: Key Features (1/2) (reuse slide 3)
# ================================================================
s = prs.slides[3]
add_text(s, 0.5, 0.3, 9, 0.5, 'Key Features (1/2)', size=28, bold=True, color=WHITE)

add_text(s, 0.5, 1.0, 4.5, 0.3, 'Multi-Source Data Ingestion', size=15, bold=True, color=ACCENT_PURPLE)
add_bullets(s, 0.5, 1.3, 4.5, 1.6, [
    'Chalk Pages: 60+ scenarios per feature, auto-scraped via Playwright',
    'Jira: Description, AC, subtasks, comments, attachments — all cached to DB',
    'LLD/HLD/Solution docs: Auto-classified (LLD/HLD/Solution/Other) with confidence',
    'All data cached in SQLite — instant on re-runs, offline-capable',
])

add_text(s, 5.5, 1.0, 4.5, 0.3, 'Intelligent Test Engine', size=15, bold=True, color=ACCENT_GREEN)
add_bullets(s, 5.5, 1.3, 4.5, 1.6, [
    'Rule-based TC generation from Chalk scenarios with enriched steps',
    'Negative scenario auto-generation (timeouts, rollbacks, auth failures)',
    'Test Analyst reasoning — thinks like a senior QA engineer',
    'Jira comment/subtask mining for additional testable scenarios',
])

add_text(s, 0.5, 3.1, 4.5, 0.3, 'Humanizer Engine', size=15, bold=True, color=ACCENT_ORANGE)
add_bullets(s, 0.5, 3.4, 4.5, 1.6, [
    'Varied descriptions — "Confirm that...", "Ensure...", "Check that..."',
    'P1/P2/P3 priority scoring based on risk analysis',
    'Smart dedup — merges TCs with >92% step overlap (title-aware)',
    'Risk-based ordering — critical scenarios first',
])

add_text(s, 5.5, 3.1, 4.5, 0.3, 'Device Matrix & Channel', size=15, bold=True, color=ACCENT_RED)
add_bullets(s, 5.5, 3.4, 4.5, 1.6, [
    'Smart Suite: ITMBO + NBOP guaranteed | eSIM + pSIM | iOS + Android',
    'Full Matrix: every combination with customizable filters',
    'CDR/Mediation features auto-detected — correct channel applied',
    'Combinations sheet in Excel — no clutter in TC descriptions',
])


# ================================================================
# SLIDE 5: Key Features (2/2) (reuse slide 4)
# ================================================================
s = prs.slides[4]
add_text(s, 0.5, 0.3, 9, 0.5, 'Key Features (2/2)', size=28, bold=True, color=WHITE)

add_text(s, 0.5, 1.0, 4.5, 0.3, 'Self-Heal Pipeline', size=15, bold=True, color=ACCENT_CYAN)
add_bullets(s, 0.5, 1.3, 4.5, 1.6, [
    '6 blocks: Jira → Chalk DB → Chalk Live → Docs → Engine → Output',
    'Each block retries up to 2x with 3-second delay on failure',
    'Failed block: "Contact Dashboard Admin — [exact error message]"',
    'Successful blocks cached — resume from failure point, not scratch',
])

add_text(s, 5.5, 1.0, 4.5, 0.3, 'Batch Mode', size=15, bold=True, color=ACCENT_GREEN)
add_bullets(s, 5.5, 1.3, 4.5, 1.6, [
    'Select multiple features — generate all suites in one click',
    'Single browser session reused across all features',
    'Per-feature download buttons with TC count in output',
    'Select All / Clear All for quick batch selection',
])

add_text(s, 0.5, 3.1, 4.5, 0.3, 'DB-Powered Traceability', size=15, bold=True, color=ACCENT_ORANGE)
add_bullets(s, 0.5, 3.4, 4.5, 1.6, [
    'Every suite saved: TCs, steps, AC, scope, warnings, engine version',
    'Jira data cached on every fetch — history + offline mode',
    'Artifact hash staleness detection — know when source data changes',
    'Cross-feature search: find TCs by keyword, category, or feature',
])

add_text(s, 5.5, 3.1, 4.5, 0.3, 'AI-Ready (LLM Layer)', size=15, bold=True, color=ACCENT_RED)
add_bullets(s, 5.5, 3.4, 4.5, 1.6, [
    'Supports: OpenAI, Azure OpenAI, AWS Bedrock, Ollama (local)',
    'Gap analysis: LLM reviews suite and suggests missing scenarios',
    'Step improvement: generic steps → feature-specific with API names',
    'Auto AI Review prompt from DB — paste into any LLM for review',
])


# ================================================================
# SLIDE 6: Excel Output (reuse slide 5)
# ================================================================
s = prs.slides[5]
add_text(s, 0.5, 0.3, 9, 0.5, 'Excel Output — Production Quality', size=28, bold=True, color=WHITE)
add_text(s, 0.5, 0.8, 9, 0.3, '4 sheets per workbook, color-coded, with Suite Guide for novice testers', size=12, color=GRAY_TEXT)

sheets_info = [
    ('Summary', ACCENT_PURPLE, ['Suite Guide: explains sheets, categories, priorities',
        'Feature details, AC, coverage breakdown, priority distribution', 'Data sources, warnings, open items']),
    ('Test Cases', ACCENT_BLUE, ['S.No, Summary, Description, Preconditions, Steps, Expected Results',
        'Category color-coding: green/red/yellow | Humanized descriptions', 'Feature-specific steps with API names and expected codes']),
    ('Traceability', ACCENT_GREEN, ['AC → TC mapping with coverage status',
        'Green = covered, Red = gap — instant visibility', 'Ensures no acceptance criteria left untested']),
    ('Combinations', ACCENT_ORANGE, ['Device matrix: Channel | Device | OS | SIM | Network',
        'ITMBO + NBOP guaranteed in Smart Suite', 'Referenced from preconditions — clean descriptions']),
]
for i, (title, color, bullets) in enumerate(sheets_info):
    top = 1.3 + i * 1.2
    add_box(s, 0.5, top, 2.5, 0.9, title, color, 12)
    add_bullets(s, 3.3, top, 7, 0.9, bullets, size=11)


# ================================================================
# SLIDE 7: Value Add (reuse slide 6)
# ================================================================
s = prs.slides[6]
add_text(s, 0.5, 0.3, 9, 0.5, 'Value Add', size=28, bold=True, color=WHITE)

values = [
    ('AI-Powered Development', 'Built entirely using Amazon Q AI — architecture to code to testing.', ACCENT_PURPLE),
    ('Zero Manual TC Writing', 'Select feature, click Execute — production-ready suite in under 2 minutes.', ACCENT_BLUE),
    ('Human-Touch Quality', 'Humanizer: varied descriptions, risk priority, smart dedup. Reads like a senior QA wrote it.', ACCENT_GREEN),
    ('Self-Healing Pipeline', '6 blocks with auto-retry. Jira slow? Retries. Chalk changed? Adapts. No manual intervention.', ACCENT_ORANGE),
    ('Complete Traceability', 'Every TC traced to AC. Every suite versioned in DB. Every artifact hashed for staleness.', ACCENT_CYAN),
    ('Reusable & Extensible', 'Modular architecture. LLM-ready. Batch mode. DB-powered. Each module independently usable.', ACCENT_RED),
]
for i, (title, desc, color) in enumerate(values):
    row, col = i // 2, i % 2
    left = 0.5 + col * 5.0
    top = 1.0 + row * 1.5
    add_text(s, left, top, 4.5, 0.3, title, size=14, bold=True, color=color)
    add_text(s, left, top + 0.35, 4.5, 0.7, desc, size=11, color=GRAY_TEXT)


# ================================================================
# SLIDE 8: Before vs After (reuse slide 7)
# ================================================================
s = prs.slides[7]
add_text(s, 0.5, 0.3, 9, 0.5, 'Benefits: Time & Resource Savings', size=28, bold=True, color=WHITE)

add_text(s, 0.5, 1.0, 4.5, 0.3, 'BEFORE (Manual Process)', size=15, bold=True, color=ACCENT_RED)
add_bullets(s, 0.5, 1.3, 4.5, 2.2, [
    'Read Chalk page manually: ~20 min per feature',
    'Read Jira story + subtasks + AC: ~15 min',
    'Write test cases in Excel: ~2-4 hours per feature',
    'Cross-check AC coverage manually: ~30 min',
    'Review and format Excel: ~30 min',
    'Total per feature: ~4-6 hours',
], size=12)

add_text(s, 5.5, 1.0, 4.5, 0.3, 'AFTER (TSG Dashboard)', size=15, bold=True, color=ACCENT_GREEN)
add_bullets(s, 5.5, 1.3, 4.5, 2.2, [
    'Auto Chalk scrape + DB cache: ~30 sec',
    'Auto Jira fetch + subtask deep-dive: ~60 sec',
    'Auto TC generation + humanization: ~5 sec',
    'Auto AC traceability + priority scoring: ~1 sec',
    'Auto Excel with 4 sheets + formatting: ~3 sec',
    'Total per feature: ~2 minutes',
], size=12)

add_text(s, 0.5, 3.8, 9, 0.3, 'Impact Summary', size=15, bold=True, color=WHITE)
add_stat(s, 0.5, 4.2, '90%', 'Time\nReduction', ACCENT_GREEN)
add_stat(s, 2.9, 4.2, '30x', 'Faster\nTC Creation', ACCENT_BLUE)
add_stat(s, 5.3, 4.2, '70+', 'TCs per\nFeature', ACCENT_PURPLE)
add_stat(s, 7.7, 4.2, 'Zero', 'Manual\nTC Writing', ACCENT_ORANGE)


# ================================================================
# SLIDE 9: Thank You (reuse slide 8)
# ================================================================
s = prs.slides[8]
add_text(s, 1.5, 2.0, 7, 0.8, 'Thank You', size=40, bold=True, color=ACCENT_PURPLE, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 3.0, 7, 0.5, 'TSG — Test Suite Generator V4.1', size=20, color=ACCENT_PURPLE, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 3.7, 7, 0.4, 'Powered by Amazon Q | Charter TMO Provisioning QA', size=14, color=GRAY_TEXT, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 4.5, 7, 0.4, 'Questions & Discussion', size=16, bold=True, color=ACCENT_PURPLE, align=PP_ALIGN.CENTER)


# ================================================================
# SAVE
# ================================================================
out = 'TSG_V4_Presentation_%s.pptx' % datetime.now().strftime('%Y%m%d')
prs.save(out)
print('Saved: %s' % out)
